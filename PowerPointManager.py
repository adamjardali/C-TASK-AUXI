from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


class PowerPointManager:
    left = Inches(-2.2)
    top = Inches(1.8)
    width = Inches(3)
    height = Inches(1.58)
    counter = 0

    def __init__(self):
        self.presentation = Presentation()

    def adjust(self):
        if PowerPointManager.counter >= 4 and PowerPointManager.counter%4 ==0:
            PowerPointManager.left = Inches(-2.3)
            PowerPointManager.top = PowerPointManager.top + Inches(1.7)
        PowerPointManager.counter += 1

    def get_left(self):
        self.adjust()
        PowerPointManager.left = PowerPointManager.left + Inches(2.3)
        return PowerPointManager.left
    
    def get_top(self):
        return PowerPointManager.top
    
    def get_width(self):
        return PowerPointManager.width
    
    def get_height(self):
        return PowerPointManager.height

    def get_slide(self, slide_number):
        counter = 0
        for slide in self.presentation.slides:
            counter += 1
            if counter == slide_number:
                return slide
        return None

    def create_slide(self):
        return self.presentation.slides.add_slide(self.presentation.slide_layouts[5])

    def get_title_shape(self,slide):
        return slide.shapes.title

    def add_title(self, slide, title_text):
        title = self.get_title_shape(slide)
        title.text = title_text
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def change_title_font(self,slide,font_name):
        title = self.get_title_shape(slide)
        title.text_frame.paragraphs[0].font.name = font_name

    def change_title_size(self,slide,font_size):
        title = self.get_title_shape(slide)
        title.text_frame.paragraphs[0].font.size = Pt(font_size)

    def change_title_color(self,slide,font_color_one,font_color_two,font_color_three):
        title = self.get_title_shape(slide)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(font_color_one,font_color_two,font_color_three)

    def change_title_to_bold(self,slide,bold):
        title = self.get_title_shape(slide)
        title.text_frame.paragraphs[0].font.bold = bold

    def add_shape(self, slide, shape_type, left, top, width, height):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        return shape


    def insert_shape_by_id(self, slide):
        # self.left = Inches(left)
        # self.top = Inches(top)
        # self.width = Inches(width)
        # self.height = Inches(height)

        left = top = width = height = Inches(5.0)


        shapes = slide.shapes
        new_shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
        # new_shape.shape_id = shape_id


    def get_shape(self,slide,shape_type):
        auto_shape_type_map = {
            'PENTAGON': MSO_SHAPE.PENTAGON,
            'CIRCLE': MSO_SHAPE.OVAL,
            'CHEVRON': MSO_SHAPE.CHEVRON
        }

        if shape_type.upper() in auto_shape_type_map:
            auto_shape_enum = auto_shape_type_map[shape_type.upper()]
            shape = slide.shapes.add_shape(auto_shape_enum, (self.get_left()), (self.get_top()), (self.get_width()), (self.get_height()))
            # left = top = width = height = Inches(1.0)

            # shape = slide.shapes.add_shape(auto_shape_enum, left,top,width,height)
            # shape.text = shape_type
            # shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
        else:
            print(f"Unsupported shape type: {shape_type}")

    def insert_auto_shape_by_type(self, slide, shape_type):
        auto_shape_type_map = {
            'PENTAGON': MSO_SHAPE.PENTAGON,
            'CIRCLE': MSO_SHAPE.OVAL,
            'CHEVRON': MSO_SHAPE.CHEVRON
        }

        if shape_type.upper() in auto_shape_type_map:
            auto_shape_enum = auto_shape_type_map[shape_type.upper()]
            shape = slide.shapes.add_shape(auto_shape_enum, (self.get_left()), (self.get_top()), (self.get_width()), (self.get_height()))
            return shape
            # left = top = width = height = Inches(1.0)

            # shape = slide.shapes.add_shape(auto_shape_enum, left,top,width,height)
            # shape.text = shape_type
            # shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
        else:
            print(f"Unsupported shape type: {shape_type}")
            return None

    def add_title_to_shape(self,shape,title):
        shape.text = title
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def get_slide_properties(self, slide):
        try:
            properties = {
                "background": slide.background,
                "follow_master_background": slide.follow_master_background,
                "has_notes_slide": slide.has_notes_slide,
                "name": slide.name,
                "placeholders": slide.placeholders,
                "shapes": slide.shapes,
                "slide_id": slide.slide_id,
                "slide_layout": slide.slide_layout.name,
                # "slide_master": slide.slide_master.name,
                "used_by_slides": slide.slide_layout.used_by_slides
            }

            print(properties)
            # Print slide dimensions
            print("Slide Dimensions:")
            # print(f"Length: {slide.height.inches} inches")
            print(f"Width: {slide.width.inches} inches")
            print(f"Left: {slide.left.inches} inches")
            print(f"Top: {slide.top.inches} inches")
            print(f"Right: {slide.left.inches + slide.width.inches} inches")
            print(f"Bottom: {slide.top.inches + slide.height.inches} inches")

            return properties
        except Exception as e:
            print(f"Error: {e}")
            return None

    def add_text_box_with_bullet_points_to_slide(self, slide, text_list):
        left = self.get_left()
        top = self.get_top()
        width = self.get_width()
        height = self.get_height()

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        
        # font = text_frame.paragraphs[0].runs[0].font
        # font.name = "Beirut"
        
        # title.text_frame.paragraphs[0].font.name = font_name

        for i in range(len(text_list)):
            p = text_frame.add_paragraph()
            p.text = "• " + text_list[i]  # Adding the bullet character manually
            # p.text.font.name = "Beirut"
            p.level = 0
            # p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            p.space_after = Pt(5)  # Adjust the spacing between bullet points
            font = p.runs[0].font
            font.name = "Beirut"
    

    

    def save_presentation(self,file_path):
        self.presentation.save(file_path)

