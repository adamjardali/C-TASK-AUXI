from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

class PowerPointAnalyzer:
    def __init__(self, file_path):
        self.presentation = Presentation(file_path)
    
    def get_slide(self, slide_number):
        counter = 0
        for slide in self.presentation.slides:
            counter += 1
            if counter == slide_number:
                return slide
        return None

    def get_title(self, slide):
        if slide.shapes.title:
            return slide.shapes.title.text
        return None

    def get_text_boxes(self, slide):
        text_boxes = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                text_boxes.append(shape.text_frame.text)
        return text_boxes
    def extract_shape_type(self, shape_type_string):
        return shape_type_string.split(' ')[0]
    def get_auto_shape_type(self, shape):
        try:
            answer = str(shape.auto_shape_type)
            output = self.extract_shape_type(answer)
            return output
        except Exception as e:
            return None
        
    def get_background_color(self, shape):
        try:
            return shape.fill.fore_color.rgb
        except Exception as e:
            return None
        
    def get_theme_color(self, shape):
        try:
            return shape.fill.fore_color.theme_color
        except Exception as e:
            return None
        
    def get_brightness(self, shape):
        try:
            return shape.fill.fore_color.brightness
        except Exception as e:
            return None
        
    def get_text_filled_in(self, shape):
        try:
            return shape.text_frame.text
        except Exception as e:
            return None
        
    def get_length(self, shape):
        try:
            return shape.height.inches
        except Exception as e:
            return None
        
    def get_width(self, shape):
        try:
            return shape.width.inches
        except Exception as e:
            return None
        
    def get_left(self, shape):
        try:
            return shape.left.inches
        except Exception as e:
            return None
        
    def get_top(self, shape):
        try:
            return shape.top.inches
        except Exception as e:
            return None
        
    def get_line_color(self, shape):
        try:
            return shape.line.color.rgb
        except Exception as e:
            return None
        
    def get_line_width(self, shape):
        try:
            return shape.line.width
        except Exception as e:
            return None


    def get_text_boxes_in_slide(self, slide):
        counter = 1
        text_boxes = []
        for shape in slide.shapes:
            if shape.has_text_frame and counter != 1 and shape.text_frame.text != '':
                text_boxes.append(shape.text_frame.text)
            counter += 1
        return text_boxes


    def get_auto_shapes(self, slide):
        auto_shapes = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # print("Auto-Shape Type:", shape.auto_shape_type)
                # iloveyou = shape.auto_shape_type
                # print(iloveyou)
                auto_shapes.append({
                    "auto_shape_type": self.get_auto_shape_type(shape),
                    "background_color": self.get_background_color(shape),
                    "theme_color": self.get_theme_color(shape),
                    "brightness": self.get_brightness(shape),
                    "text_filled_in": self.get_text_filled_in(shape),
                    "length": self.get_length(shape),
                    "width": self.get_width(shape),
                    "left": self.get_left(shape),
                    "top": self.get_top(shape),
                    "line_color": self.get_line_color(shape),
                    "line_width": self.get_line_width(shape)
                })
        # print(auto_shapes)
        return auto_shapes

    def analyze_presentation(self):
        for slide in self.presentation.slides:
            # print("Slide Number:", slide.slide_id)
            # title = self.get_title(slide)
            # if title:
            #     print("Title:", title)
            
            # text_boxes = self.get_text_boxes(slide)
            # if text_boxes:
            #     for i, text_box in enumerate(text_boxes):
            #         print(f"Text Box {i+1}:", text_box)

            auto_shapes = self.get_auto_shapes(slide)
            if auto_shapes:
                for i, shape_data in enumerate(auto_shapes):
                    print(f"Auto-Shape {i+1} Type:", shape_data["auto_shape_type"])
                    if shape_data["background_color"]:
                        print("Background Color:", shape_data["background_color"])
                    if shape_data["text_filled_in"]:
                        print("Text Filled In:", shape_data["text_filled_in"])
                    print("Length:", shape_data["length"])
                    print("Width:", shape_data["width"])
